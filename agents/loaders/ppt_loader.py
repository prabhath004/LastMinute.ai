from pptx import Presentation

from .base_loader import BaseLoader


class PPTLoader(BaseLoader):
    def load(self, path: str) -> str:
        texts = []
        presentation = Presentation(path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
